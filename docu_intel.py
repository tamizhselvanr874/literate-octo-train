import json  
import streamlit as st  
from pptx import Presentation  
from pptx.enum.shapes import MSO_SHAPE_TYPE  
from PIL import Image  
from io import BytesIO  
import requests  
import base64  
from docx import Document  
from docx.shared import Pt  
import fitz  # PyMuPDF  
import os  
import cv2  
import numpy as np 
  
# Azure OpenAI credentials  
azure_endpoint = "https://gpt-4omniwithimages.openai.azure.com/"  
api_key = "6e98566acaf24997baa39039b6e6d183"  
api_version = "2024-02-01"  
model = "GPT-40-mini"  
  
# Azure Graph API credentials  
GRAPH_TENANT_ID = "4d4343c6-067a-4794-91f3-5cb10073e5b4"  
GRAPH_CLIENT_ID = "5ace14db-3235-4cd2-acfd-dd5ef19d6ea1"  
GRAPH_CLIENT_SECRET = "HRk8Q~7G6EH3.yhDC3rB5wLAyAixQMnQNWNyUdsW"  
PDF_SITE_ID = "marketingai.sharepoint.com,b82dbaac-09cc-4539-ad08-e4ca926796e8,7b756d20-3463-44b7-95ca-5873f8c3f517"  
  
# Function to get OAuth2 token  
def get_oauth2_token():  
    url = f"https://login.microsoftonline.com/{GRAPH_TENANT_ID}/oauth2/v2.0/token"  
    headers = {  
        'Content-Type': 'application/x-www-form-urlencoded'  
    }  
    data = {  
        'grant_type': 'client_credentials',  
        'client_id': GRAPH_CLIENT_ID,  
        'client_secret': GRAPH_CLIENT_SECRET,  
        'scope': 'https://graph.microsoft.com/.default'  
    }  
    response = requests.post(url, headers=headers, data=data)  
    if response.status_code == 200:  
        return response.json().get('access_token')  
    else:  
        st.error(f"Failed to obtain OAuth2 token: {response.content}")  
        return None  
  
# Function to upload file to SharePoint  
def upload_file_to_sharepoint(token, file_path):  
    with open(file_path, "rb") as file:  
        upload_url = f"https://graph.microsoft.com/v1.0/sites/{PDF_SITE_ID}/drive/root:/{os.path.basename(file_path)}:/content"  
        headers = {  
            'Authorization': f'Bearer {token}',  
            'Content-Type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'  
        }  
        response = requests.put(upload_url, headers=headers, data=file.read())  
        if response.status_code == 201:  
            return response.json().get('id')  
        else:  
            st.error(f"Failed to upload file to SharePoint: {response.content}")  
            return None  
  
# Function to convert file to PDF using Microsoft Graph API  
def convert_file_to_pdf(token, file_id):  
    convert_url = f"https://graph.microsoft.com/v1.0/sites/{PDF_SITE_ID}/drive/items/{file_id}/content?format=pdf"  
    headers = {  
        'Authorization': f'Bearer {token}',  
        'Content-Type': 'application/json'  
    }  
    response = requests.get(convert_url, headers=headers)  
    if response.status_code == 200:  
        return response.content  
    else:  
        st.error(f"Failed to convert file to PDF: {response.content}")  
        return None  
  
# Function to delete file from SharePoint  
def delete_file_from_sharepoint(token, file_id):  
    delete_url = f"https://graph.microsoft.com/v1.0/sites/{PDF_SITE_ID}/drive/items/{file_id}"  
    headers = {  
        'Authorization': f'Bearer {token}'  
    }  
    response = requests.delete(delete_url, headers=headers)  
    if response.status_code == 204:  
        return True  
    else:  
        st.error(f"Failed to delete file from SharePoint: {response.content}")  
        return False  
  
# Function to encode image as base64  
def encode_image(image):  
    return base64.b64encode(image).decode("utf-8")  
  
def get_image_explanation(base64_image):  
    headers = {  
        "Content-Type": "application/json",  
        "api-key": api_key  
    }  
    data = {  
        "model": model,  
        "messages": [  
            {"role": "system", "content": "You are a helpful assistant that responds in Markdown."},  
            {"role": "user", "content": [  
                {"type": "text", "text": "Explain the content of this image in a single, coherent paragraph. The explanation should be concise and semantically meaningful, summarizing all major points from the image in one continuous paragraph. Avoid using bullet points, line breaks, or separate lists."},  
                {"type": "image_url", "image_url": {  
                    "url": f"data:image/png;base64,{base64_image}"}  
                }  
            ]}  
        ],  
        "temperature": 0.7  
    }  
  
    response = requests.post(  
        f"{azure_endpoint}/openai/deployments/{model}/chat/completions?api-version={api_version}",  
        headers=headers,  
        json=data  
    )  
  
    if response.status_code == 200:  
        result = response.json()  
        return result["choices"][0]["message"]["content"]  
    else:  
        st.error(f"Error: {response.status_code} - {response.text}")  
        return None  
  
def ppt_to_pdf(ppt_file, pdf_file):  
    token = get_oauth2_token()  
    if token:  
        file_id = upload_file_to_sharepoint(token, ppt_file)  
        if file_id:  
            pdf_content = convert_file_to_pdf(token, file_id)  
            if pdf_content:  
                with open(pdf_file, "wb") as pdf_file:  
                    pdf_file.write(pdf_content)  
                delete_file_from_sharepoint(token, file_id)  
                return True  
    return False  
  
def extract_text_from_ppt(ppt_file):  
    presentation = Presentation(ppt_file)  
    text_content = []  
    for slide_number, slide in enumerate(presentation.slides, start=1):  
        slide_text = []  
        for shape in slide.shapes:  
            if hasattr(shape, "text"):  
                slide_text.append(shape.text)  
        slide_title = slide.shapes.title.text if slide.shapes.title else "Untitled Slide"  
        text_content.append({"slide_number": slide_number, "slide_title": slide_title, "text": " ".join(slide_text)})  
    return text_content  
  
def identify_visual_elements(ppt_file):  
    presentation = Presentation(ppt_file)  
    visual_slides = []  
    for slide_number, slide in enumerate(presentation.slides, start=1):  
        has_visual_elements = False  
        for shape in slide.shapes:  
            if shape.shape_type in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.TABLE, MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.AUTO_SHAPE}:  
                has_visual_elements = True  
                break  
        if has_visual_elements:  
            visual_slides.append(slide_number)  
    return visual_slides  
  
def capture_slide_images(pdf_file, slide_numbers):  
    doc = fitz.open(pdf_file)  
    images = []  
    for slide_number in slide_numbers:  
        page = doc[slide_number - 1]  
        pix = page.get_pixmap()  
        image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)  
        buffer = BytesIO()  
        image.save(buffer, format="PNG")  
        images.append({"slide_number": slide_number, "image": buffer.getvalue()})  
    return images  
  
def generate_text_insights(text_content, visual_slides, text_length):  
    headers = {  
        "Content-Type": "application/json",  
        "api-key": api_key  
    }  
    insights = []  
  
    # Set temperature based on text_length  
    if text_length == "Standard":  
        temperature = 0.3  
    elif text_length == "Blend":  
        temperature = 0.5  
    elif text_length == "Creative":  
        temperature = 0.7  
  
    for slide in text_content:  
        slide_text = slide['text']  
        slide_number = slide['slide_number']  
        if len(slide_text.split()) < 20 and slide_number not in visual_slides:  
            continue  # Skip slides with fewer than 20 words and no visual elements  
        prompt = f"""  
        Aspects of the present disclosure may include insights extracted from the above slide content. The information should be delivered directly and engagingly in a single, coherent paragraph. Avoid phrases like 'The slide presents,' 'discusses,' 'outlines,' or 'content.' The explanation should be concise and semantically meaningful, summarizing all major points in one paragraph without line breaks or bullet points. The text should adhere to the following style guidelines:  
        1. Remove all listed profanity words.  
        2. Use passive voice.  
        3. Use conditional and tentative language, such as "may include," "in some aspects," and "aspects of the present disclosure."  
        4. Replace "Million" with "1,000,000" and "Billion" with "1,000,000,000."  
        5. Maintain the following tone characteristics: Precision and Specificity, Formality, Complexity, Objective and Impersonal, Structured and Systematic.  
        6. Follow these style elements: Formal and Objective, Structured and Systematic, Technical Jargon and Terminology, Detailed and Specific, Impersonal Tone, Instructional and Descriptive, Use of Figures and Flowcharts, Legal and Protective Language, Repetitive and Redundant, Examples and Clauses.  
        7. Use the following conditional and tentative language phrases: may include, in some aspects, aspects of the present disclosure, wireless communication networks, by way of example, may be, may further include, may be used, may occur, may use, may monitor, may periodically wake up, may demodulate, may consume, can be performed, may enter and remain, may correspond to, may also include, may be identified in response to, may be further a function of, may be multiplied by, may schedule, may select, may also double, may further comprise, may be configured to, may correspond to a duration value, may correspond to a product of, may be closer, may be significant, may not be able, may result, may reduce, may be operating in, may further be configured to, may further process, may be executed by, may be received, may avoid, may indicate, may be selected, may be proactive, may perform, may be necessary, may be amplified, may involve, may require, may be stored, may be accessed, may be transferred, may be implemented, may include instructions to, may depend upon, may communicate, may be generated, may be configured.  
        8. Maintain the exact wording in the generated content. Do not substitute words with synonyms. For example, "instead" should remain "instead" and not be replaced with "conversely."  
        9. Replace the phrase "further development" with "our disclosure" in all generated content.  
        10. Make sure to use LaTeX formatting for all mathematical symbols, equations, subscripting, and superscripting to ensure they are displayed correctly in the output.  
        11. When encountering programmatic terms or equations, ensure they are accurately represented and contextually retained.
        12. Avoid using the term consist or any form of that verb when describing the invention.
        13. For slides with the heading including Background, treat the content as a prior solution and not as our disclosure.
        14. Ensure that the term antennas is used consistently throughout the paragraph after its initial mention.
        15. Ensure all information stated is accurate and contextually correct.
        16. Summarize the content effectively, ensuring the meaning of the wording is retained.
        17. Capture all key wording and phrases accurately.
        18. Reference figures accurately and appropriately within the context.
        19. Avoid repeatedly spelling out abbreviations if they have already been defined above.
        20. Place the summary of the slide at the beginning if it is the first bullet point.
        21. For slides with Background in the heading, avoid highlighting any advantages or efficiencies.
        22. When discussing our disclosure, use definitive language.
        23. Avoid duplicating descriptions; ensure the content is concise and not repetitive.  
        {slide_text}  
        """  
        if text_length == "Standard":  
            prompt += "\n\nGenerate a short paragraph."  
        elif text_length == "Blend":  
            prompt += "\n\nGenerate a medium-length paragraph."  
        elif text_length == "Creative":  
            prompt += "\n\nGenerate a longer paragraph."  
  
        data = {  
            "model": model,  
            "messages": [{"role": "system", "content": "You are a helpful assistant."}, {"role": "user", "content": prompt}],  
            "temperature": temperature  
        }  
  
        response = requests.post(  
            f"{azure_endpoint}/openai/deployments/{model}/chat/completions?api-version={api_version}",  
            headers=headers,  
            json=data  
        )  
  
        if response.status_code == 200:  
            result = response.json()  
            insights.append({"slide_number": slide['slide_number'], "slide_title": slide['slide_title'], "insight": result["choices"][0]["message"]["content"]})  
        else:  
            st.error(f"Error: {response.status_code} - {response.text}")  
            insights.append({"slide_number": slide['slide_number'], "slide_title": slide['slide_title'], "insight": "Error in generating insight"})  
  
    return insights  
  
def generate_image_insights(image_content):  
    insights = []  
    temperature = 0.5  # Fixed temperature  
  
    for image_data in image_content:  
        base64_image = encode_image(image_data['image'])  
        headers = {  
            "Content-Type": "application/json",  
            "api-key": api_key  
        }  
        prompt = f"""  
            Analyze and explain the content of this image by first analyzing the text and then describing the image. Reproduce the text as accurately as possible, interweaving your explanation of the image with the text to maintain context. If the image contains multiple figures, reference each one separately while ensuring the explanation flows smoothly. Provide a concise and semantically meaningful summary of the entire image in a single, coherent paragraph, covering all major points without using bullet points, line breaks, or separate lists. Ensure that the explanation is both detailed and continuous, blending the textual and visual information effectively.  
  
            The explanation should adhere to the following style guidelines:  
            1. Remove all listed profanity words.  
            2. Use passive voice.  
            3. Use conditional and tentative language, such as "may include," "in some aspects," and "aspects of the present disclosure."  
            4. Replace "Million" with "1,000,000" and "Billion" with "1,000,000,000."  
            5. Maintain the following tone characteristics: Precision and Specificity, Formality, Complexity, Objective and Impersonal, Structured and Systematic.  
            6. Follow these style elements: Formal and Objective, Structured and Systematic, Technical Jargon and Terminology, Detailed and Specific, Impersonal Tone, Instructional and Descriptive, Use of Figures and Flowcharts, Legal and Protective Language, Repetitive and Redundant, Examples and Clauses.  
            7. Use the following conditional and tentative language phrases: may include, in some aspects, aspects of the present disclosure, wireless communication networks, by way of example, may be, may further include, may be used, may occur, may use, may monitor, may periodically wake up, may demodulate, may consume, can be performed, may enter and remain, may correspond to, may also include, may be identified in response to, may be further a function of, may be multiplied by, may schedule, may select, may also double, may further comprise, may be configured to, may correspond to a duration value, may correspond to a product of, may be closer, may be significant, may not be able, may result, may reduce, may be operating in, may further be configured to, may further process, may be executed by, may be received, may avoid, may indicate, may be selected, may be proactive, may perform, may be necessary, may be amplified, may involve, may require, may be stored, may be accessed, may be transferred, may be implemented, may include instructions to, may depend upon, may communicate, may be generated, may be configured.  
            8. Maintain the exact wording in the generated content. Do not substitute words with synonyms. For example, "instead" should remain "instead" and not be replaced with "conversely."  
            9. Replace the phrase "further development" with "our disclosure" in all generated content.  
            10. Make sure to use LaTeX formatting for all mathematical symbols, equations, subscripting, and superscripting to ensure they are displayed correctly in the output.  
            11. When encountering programmatic terms or equations, ensure they are accurately represented and contextually retained.  
            12. Avoid using the term consist or any form of that verb when describing the invention.  
            13. For slides with the heading including Background, treat the content as a prior solution and not as our disclosure.  
            14. Ensure that the term antennas is used consistently throughout the paragraph after its initial mention.  
            15. Ensure all information stated is accurate and contextually correct.  
            16. Summarize the content effectively, ensuring the meaning of the wording is retained.  
            17. Capture all key wording and phrases accurately.  
            18. Reference figures accurately and appropriately within the context.  
            19. Avoid repeatedly spelling out abbreviations if they have already been defined above.  
            20. Place the summary of the slide at the beginning if it is the first bullet point.  
            21. For slides with Background in the heading, avoid highlighting any advantages or efficiencies.  
            22. When discussing our disclosure, use definitive language.  
            23. Avoid duplicating descriptions; ensure the content is concise and not repetitive.  
        """  
  
        data = {  
            "model": model,  
            "messages": [  
                {"role": "system", "content": "You are a helpful assistant that responds in Markdown."},  
                {"role": "user", "content": prompt},  
                {"role": "user", "content": {  
                    "type": "image_url",  
                    "image_url": {  
                        "url": f"data:image/png;base64,{base64_image}"  
                    }  
                }}  
            ],  
            "temperature": temperature  
        }  
  
        response = requests.post(  
            f"{azure_endpoint}/openai/deployments/{model}/chat/completions?api-version={api_version}",  
            headers=headers,  
            json=data  
        )  
  
        if response.status_code == 200:  
            result = response.json()  
            insights.append({"slide_number": image_data['slide_number'], "slide_title": image_data.get('slide_title', 'Untitled Slide'), "insight": result["choices"][0]["message"]["content"]})  
        else:  
            st.error(f"Error: {response.status_code} - {response.text}")  
            insights.append({"slide_number": image_data['slide_number'], "slide_title": image_data.get('slide_title', 'Untitled Slide'), "insight": "Error in generating insight"})  
  
    return insights  
 
  
def aggregate_content(text_insights, image_insights):  
    aggregated_content = []  
    for text in text_insights:  
        slide_number = text['slide_number']  
        slide_title = text['slide_title']  
        text_insight = text['insight']  
        image_insight = next((img['insight'] for img in image_insights if img['slide_number'] == slide_number), None)  
        if image_insight:  
            aggregated_content.append({  
                "slide_number": slide_number,  
                "slide_title": slide_title,  
                "content": f"Referring to Figure {slide_number}. {image_insight} {text_insight}"  
            })  
        else:  
            aggregated_content.append({  
                "slide_number": slide_number,  
                "slide_title": slide_title,  
                "content": text_insight  
            })  
    return aggregated_content  
  
def sanitize_text(text):  
    if text:  
        sanitized = ''.join(c for c in text if c.isprintable() and c not in {'\x00', '\x01', '\x02', '\x03', '\x04', '\x05', '\x06', '\x07', '\x08', '\x0B', '\x0C', '\x0E', '\x0F', '\x10', '\x11', '\x12', '\x13', '\x14', '\x15', '\x16', '\x17', '\x18', '\x19', '\x1A', '\x1B', '\x1C', '\x1D', '\x1E', '\x1F'})  
        return sanitized  
    return text  
  
def save_content_to_word(aggregated_content, output_file_name, extracted_images):  
    doc = Document()  
    style = doc.styles['Normal']  
    font = style.font  
    font.name = 'Times New Roman'  
    font.size = Pt(10.5)  # Reduced font size for paragraphs  
    paragraph_format = style.paragraph_format  
    paragraph_format.line_spacing = 1.5  
    paragraph_format.alignment = 3  # Justify  
  
    for slide in aggregated_content:  
        sanitized_title = sanitize_text(slide['slide_title'])  
        sanitized_content = sanitize_text(slide['content'])  
        doc.add_heading(f"[[{slide['slide_number']}, {sanitized_title}]]", level=1)  
        if sanitized_content:  # Only add content if it exists  
            doc.add_paragraph(sanitized_content)  
  
    # Add extracted images after the generated content  
    if extracted_images:  
        doc.add_heading("Extracted Images", level=1)  
        for idx, (image, slide_number) in enumerate(extracted_images):  
            _, buffer = cv2.imencode('.png', image)  
            image_stream = BytesIO(buffer)  
            doc.add_paragraph(f"Image from Slide {slide_number}:")  
            doc.add_picture(image_stream, width=doc.sections[0].page_width - doc.sections[0].left_margin - doc.sections[0].right_margin)  
            doc.add_paragraph("\n")  # Add space after image  
  
    output = BytesIO()  
    doc.save(output)  
    output.seek(0)  
    return output  
  
def extract_and_clean_page_image(page, top_mask, bottom_mask, left_mask, right_mask):  
    # Get the page as an image  
    pix = page.get_pixmap()  
    img_data = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, pix.n)  
  
    # Convert the image to BGR format for OpenCV  
    img_bgr = cv2.cvtColor(img_data, cv2.COLOR_RGB2BGR)  
  
    # Convert to grayscale for processing  
    gray = cv2.cvtColor(img_bgr, cv2.COLOR_BGR2GRAY)  
  
    # Threshold the image to get binary image  
    _, binary = cv2.threshold(gray, 240, 255, cv2.THRESH_BINARY_INV)  
  
    # Detect contours of possible images/diagrams  
    contours, _ = cv2.findContours(binary, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)  
  
    # Check if there are any valid contours (image regions)  
    valid_contours = [cv2.boundingRect(contour) for contour in contours if cv2.boundingRect(contour)[2] > 50 and cv2.boundingRect(contour)[3] > 50]  
    if not valid_contours:  
        return None  # Skip the page if no valid images/diagrams are found  
  
    # Create a mask for the detected contours
    mask = np.zeros_like(binary)  
    for x, y, w, h in valid_contours:  
        # Apply the adjustable top, bottom, left, and right masking values from the sliders  
        # Ensure coordinates do not go out of image bounds  
        x1 = max(0, x - left_mask)  
        y1 = max(0, y - top_mask)  
        x2 = min(img_bgr.shape[1], x + w + right_mask)  
        y2 = min(img_bgr.shape[0], y + h + bottom_mask)  
        cv2.rectangle(mask, (x1, y1), (x2, y2), 255, -1)  
  
    # Use the mask to keep only the regions with images/diagrams  
    text_removed = cv2.bitwise_and(img_bgr, img_bgr, mask=mask)  
  
    # Set the background to white where the mask is not applied  
    background = np.ones_like(img_bgr) * 255  
    cleaned_image = np.where(mask[:, :, None] == 255, text_removed, background)  
  
    # Convert cleaned image to grayscale  
    cleaned_image_gray = cv2.cvtColor(cleaned_image, cv2.COLOR_BGR2GRAY)  
    return cleaned_image_gray  
  
def extract_images_from_pdf(pdf_file, top_mask, bottom_mask, left_mask, right_mask):  
    # Open the PDF file  
    pdf_document = fitz.open(pdf_file)  
    page_images = []  
  
    for page_num in range(len(pdf_document)):  
        page = pdf_document.load_page(page_num)  
  
        # Extract and clean the page image  
        cleaned_image = extract_and_clean_page_image(page, top_mask, bottom_mask, left_mask, right_mask)  
        if cleaned_image is not None:  
            page_images.append((cleaned_image, page_num + 1))  # Keep track of the slide number  
  
    pdf_document.close()  
    return page_images  

def main():  
    st.title("PPT Insights Extractor")   
  
    text_length = st.select_slider(  
        "Content Generation Slider",  
        options=["Standard", "Blend", "Creative"],  
        value="Blend"  
    )  
  
    # Add Title and Information Button for Image Extraction Slider  
    st.sidebar.markdown("### Image Extraction Slider")  
  
    # Initialize session state variables for the sliders  
    if 'top_mask' not in st.session_state:  
        st.session_state.top_mask = 40  
    if 'bottom_mask' not in st.session_state:  
        st.session_state.bottom_mask = 40  
    if 'left_mask' not in st.session_state:  
        st.session_state.left_mask = 85  
    if 'right_mask' not in st.session_state:  
        st.session_state.right_mask = 85  
  
    # Arrange the buttons in a row using columns  
    col1, col2 = st.sidebar.columns(2)  
    with col1:  
        if st.button("Default"):  
            st.session_state.top_mask = 40  
            st.session_state.bottom_mask = 40  
            st.session_state.left_mask = 85  
            st.session_state.right_mask = 85  
  
    with col2:  
        if st.button("A4"):  
            st.session_state.top_mask = 70  
            st.session_state.bottom_mask = 70  
            st.session_state.left_mask = 85  
            st.session_state.right_mask = 85  
  
    # Add sliders to adjust the top, bottom, left, and right masking values  
    top_mask = st.sidebar.slider("Adjust Top Masking Value", min_value=10, max_value=100, value=st.session_state.top_mask, step=1)  
    bottom_mask = st.sidebar.slider("Adjust Bottom Masking Value", min_value=10, max_value=100, value=st.session_state.bottom_mask, step=1)  
    left_mask = st.sidebar.slider("Adjust Left Masking Value", min_value=10, max_value=500, value=st.session_state.left_mask, step=1)  
    right_mask = st.sidebar.slider("Adjust Right Masking Value", min_value=10, max_value=200, value=st.session_state.right_mask, step=1)  
  
    # Update session state if sliders are moved  
    if top_mask != st.session_state.top_mask or bottom_mask != st.session_state.bottom_mask or left_mask != st.session_state.left_mask or right_mask != st.session_state.right_mask:  
        st.session_state.top_mask = top_mask  
        st.session_state.bottom_mask = bottom_mask  
        st.session_state.left_mask = left_mask  
        st.session_state.right_mask = right_mask  
  
    uploaded_ppt = st.file_uploader("Upload a PPT file", type=["pptx"])  
  
    if uploaded_ppt is not None:  
        st.info("Processing PPT file...")  
  
        # Extract the base name of the uploaded PPT file  
        ppt_filename = uploaded_ppt.name  
        base_filename = os.path.splitext(ppt_filename)[0]  
        output_word_filename = f"{base_filename}.docx"  
  
        try:  
            # Convert PPT to PDF  
            with open("temp_ppt.pptx", "wb") as f:  
                f.write(uploaded_ppt.read())  
            if not ppt_to_pdf("temp_ppt.pptx", "temp_pdf.pdf"):  
                st.error("PDF conversion failed. Please check the uploaded PPT file.")  
                return  
  
            # Extract text and identify slides with visual elements  
            text_content = extract_text_from_ppt("temp_ppt.pptx")  
            visual_slides = identify_visual_elements("temp_ppt.pptx")  
  
            # Capture images of marked slides  
            slide_images = capture_slide_images("temp_pdf.pdf", visual_slides)  
  
            st.info("Generating text insights...")  
            text_insights = generate_text_insights(text_content, visual_slides, text_length)  
  
            st.info("Generating image insights...")  
            image_insights = generate_image_insights(slide_images, text_length)  
  
            st.info("Extracting additional images...")  
            extracted_images = extract_images_from_pdf("temp_pdf.pdf", top_mask, bottom_mask, left_mask, right_mask)  
  
            st.info("Aggregating content...")  
            aggregated_content = aggregate_content(text_insights, image_insights)  
  
            st.info("Saving to Word document...")  
            output_doc = save_content_to_word(aggregated_content, output_word_filename, extracted_images)  
  
            st.download_button(label="Download Word Document", data=output_doc, file_name=output_word_filename)  
  
            st.success("Processing completed successfully!")  
        except Exception as e:  
            st.error(f"An error occurred: {e}")  
  
if __name__ == "__main__":  
    main()  

